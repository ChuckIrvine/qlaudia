#!/usr/bin/env python3
"""Convert the Qlaudia pitch deck markdown into a PowerPoint file.

Structure understood by this converter (see the markdown's own "Convention" block):

    # Title                     -> title slide (leading blockquote becomes its notes)
    ## Slide N - Title          -> one content slide
    ## Appendix ...             -> section divider; following ### become their own slides
    **Headline:** ...           -> the large assertion under the slide title
    - / 1.                      -> bullets
    | a | b |                   -> native PowerPoint table
    > **Speaker notes:** ...    -> the notes pane, NOT the slide
    > ...                       -> pull quote rendered on the slide
    **[VISUAL: ...]**           -> a dashed image placeholder

Usage:  .venv/bin/python md2pptx.py [source.md] [output.pptx]
"""

from __future__ import annotations

import re
import sys
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------- design tokens

SLIDE_W, SLIDE_H = 13.333, 7.5

INK = RGBColor(0x15, 0x19, 0x2B)
MUTED = RGBColor(0x5A, 0x60, 0x72)
ACCENT = RGBColor(0x4A, 0x3A, 0xFF)
TODO = RGBColor(0xC2, 0x41, 0x0C)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
WASH = RGBColor(0xF4, 0xF5, 0xF9)
RULE = RGBColor(0xD8, 0xDB, 0xE6)
DARK = RGBColor(0x15, 0x19, 0x2B)
DARK_MUTED = RGBColor(0x9A, 0xA1, 0xB8)

M_L, M_R, M_T, M_B = 0.62, 0.62, 0.46, 0.42
BODY_W = SLIDE_W - M_L - M_R
COL_GAP = 0.42

FONT = "Segoe UI"          # falls back gracefully on macOS PowerPoint
FONT_MONO = "Consolas"

TODO_MARK = "⚠"        # the warning sign used in the markdown


# ---------------------------------------------------------------- block model

@dataclass
class Block:
    kind: str                      # subhead | bullet | number | para | quote | table | visual
    text: str = ""
    level: int = 0
    rows: list[list[str]] = field(default_factory=list)


@dataclass
class Slide:
    title: str = ""
    headline: str = ""
    blocks: list[Block] = field(default_factory=list)
    notes: list[str] = field(default_factory=list)
    kind: str = "content"          # content | title | divider


# ---------------------------------------------------------------- markdown parse

INLINE_RE = re.compile(r"(\*\*.+?\*\*|\*[^*]+?\*|`[^`]+?`)", re.S)
LINK_RE = re.compile(r"\[([^\]]+)\]\(([^)]+)\)")


def parse(md: str) -> list[Slide]:
    lines = md.splitlines()
    slides: list[Slide] = []
    cur: Slide | None = None
    in_appendix = False
    pending_notes = False            # inside a "> **Speaker notes:**" run
    i = 0

    def flush_para(buf: list[str]) -> None:
        text = " ".join(x.strip() for x in buf).strip()
        if text and cur is not None:
            cur.blocks.append(Block("para", text))
        buf.clear()

    para: list[str] = []

    while i < len(lines):
        raw = lines[i]
        line = raw.rstrip()
        stripped = line.strip()

        # --- horizontal rules and blanks -------------------------------------
        if stripped in ("---", "***"):
            flush_para(para)
            pending_notes = False
            i += 1
            continue
        if not stripped:
            flush_para(para)
            pending_notes = False
            i += 1
            continue

        # --- headings ---------------------------------------------------------
        if stripped.startswith("# "):
            flush_para(para)
            cur = Slide(title=stripped[2:].strip(), kind="title")
            slides.append(cur)
            pending_notes = False
            i += 1
            continue

        if stripped.startswith("## "):
            flush_para(para)
            title = stripped[3:].strip()
            if title.lower().startswith("appendix"):
                in_appendix = True
                cur = Slide(title=title, kind="divider")
            else:
                cur = Slide(title=title)
            slides.append(cur)
            pending_notes = False
            i += 1
            continue

        if stripped.startswith("### "):
            flush_para(para)
            text = stripped[4:].strip()
            if in_appendix:
                cur = Slide(title=text)
                slides.append(cur)
            elif cur is not None:
                cur.blocks.append(Block("subhead", text))
            pending_notes = False
            i += 1
            continue

        if cur is None:
            i += 1
            continue

        # --- blockquotes: speaker notes or pull quote -------------------------
        if stripped.startswith(">"):
            flush_para(para)
            body = stripped[1:].strip()
            plain = strip_marks(body)

            if plain.lower().startswith("speaker notes"):
                pending_notes = True
                rest = body.split(":", 1)[1].strip() if ":" in body else ""
                if strip_marks(rest):
                    cur.notes.append(rest)
                i += 1
                continue

            if not body:                      # the ">" spacer between note paras
                i += 1
                continue

            # A title slide's leading blockquote is deck metadata -> notes.
            if cur.kind == "title":
                cur.notes.append(body)
            elif pending_notes:
                cur.notes.append(body)
            else:
                cur.blocks.append(Block("quote", body))
            i += 1
            continue

        pending_notes = False

        # --- tables -----------------------------------------------------------
        if stripped.startswith("|"):
            flush_para(para)
            rows: list[list[str]] = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                cells = [c.strip() for c in lines[i].strip().strip("|").split("|")]
                if not all(re.fullmatch(r":?-{2,}:?", c or "-") for c in cells):
                    rows.append(cells)
                i += 1
            if rows:
                cur.blocks.append(Block("table", rows=rows))
            continue

        # --- headline ---------------------------------------------------------
        if stripped.startswith("**Headline:**"):
            flush_para(para)
            cur.headline = stripped[len("**Headline:**"):].strip()
            i += 1
            continue

        # --- visual placeholder ----------------------------------------------
        if stripped.startswith("**[VISUAL:") or stripped.startswith("[VISUAL:"):
            flush_para(para)
            inner = stripped.strip("*").strip()[1:-1]
            if inner.lower().startswith("visual:"):
                inner = inner.split(":", 1)[1].strip()
            cur.blocks.append(Block("visual", inner))
            i += 1
            continue

        # --- list items -------------------------------------------------------
        m = re.match(r"^(\s*)([-*])\s+(.*)$", raw)
        if m:
            flush_para(para)
            level = len(m.group(1)) // 2
            cur.blocks.append(Block("bullet", m.group(3).strip(), level=level))
            i += 1
            continue

        m = re.match(r"^(\s*)(\d+)\.\s+(.*)$", raw)
        if m:
            flush_para(para)
            level = len(m.group(1)) // 2
            cur.blocks.append(Block("number", f"{m.group(2)}. {m.group(3).strip()}", level=level))
            i += 1
            continue

        para.append(stripped)
        i += 1

    flush_para(para)
    return slides


def strip_marks(text: str) -> str:
    text = LINK_RE.sub(r"\1", text)
    return text.replace("**", "").replace("`", "").replace("*", "").strip()


# ---------------------------------------------------------------- text runs

def add_runs(paragraph, text: str, size: float, color: RGBColor, bold_color: RGBColor | None = None):
    """Render inline **bold**, *italic* and `code` as real PowerPoint runs."""
    text = LINK_RE.sub(r"\1", text)
    highlight = TODO_MARK in text

    for piece in INLINE_RE.split(text):
        if not piece:
            continue
        bold = italic = mono = False
        if piece.startswith("**") and piece.endswith("**") and len(piece) > 4:
            piece, bold = piece[2:-2], True
        elif piece.startswith("*") and piece.endswith("*") and len(piece) > 2:
            piece, italic = piece[1:-1], True
        elif piece.startswith("`") and piece.endswith("`") and len(piece) > 2:
            piece, mono = piece[1:-1], True

        run = paragraph.add_run()
        run.text = piece
        f = run.font
        f.size = Pt(size)
        f.name = FONT_MONO if mono else FONT
        f.bold = bold
        f.italic = italic
        if highlight:
            f.color.rgb = TODO
        elif bold and bold_color is not None:
            f.color.rgb = bold_color
        else:
            f.color.rgb = color
    return paragraph


# ---------------------------------------------------------------- measurement

def est_lines(text: str, size: float, width_in: float) -> int:
    per_char = size * 0.0080
    cpl = max(12, int(width_in / per_char))
    n = max(1, len(strip_marks(text)))
    return max(1, -(-n // cpl))


BODY_PT = 13.5
SUBHEAD_PT = 13.0
QUOTE_PT = 13.0
TABLE_PT = 11.0


def block_height(b: Block, width_in: float, scale: float) -> float:
    if b.kind == "table":
        head = 0.34 * scale
        return head + 0.315 * scale * (len(b.rows) - 1) + 0.20
    if b.kind == "subhead":
        return 0.30 * scale + 0.10
    if b.kind == "visual":
        return 1.55
    size = QUOTE_PT if b.kind == "quote" else BODY_PT
    size *= scale
    indent = 0.30 * b.level + (0.34 if b.kind == "quote" else 0.0)
    lines = est_lines(b.text, size, width_in - indent - (0.26 if b.kind in ("bullet",) else 0.0))
    gap = 0.11 if b.kind in ("bullet", "number") else 0.16
    return lines * (size * 1.30 / 72.0) + gap


# ---------------------------------------------------------------- rendering

def txbox(slide, l, t, w, h):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return box, tf


def bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def hline(slide, l, t, w, color=RULE, thick=0.014):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(thick))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def render_title(prs, s: Slide):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, DARK)

    _, tf = txbox(slide, M_L + 0.3, 2.35, BODY_W - 1.2, 2.2)
    p = tf.paragraphs[0]
    main = s.title.split("—")[0].strip()
    sub = s.title.split("—", 1)[1].strip() if "—" in s.title else ""
    r = p.add_run()
    r.text = main
    r.font.size, r.font.bold, r.font.name = Pt(60), True, FONT
    r.font.color.rgb = PAPER

    hline(slide, M_L + 0.3, 3.55, 1.6, ACCENT, 0.035)

    if sub:
        _, tf2 = txbox(slide, M_L + 0.3, 3.85, BODY_W - 1.2, 0.9)
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = sub
        r2.font.size, r2.font.name = Pt(22), FONT
        r2.font.color.rgb = DARK_MUTED

    add_notes(slide, s.notes)
    return slide


def render_divider(prs, s: Slide):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, DARK)
    _, tf = txbox(slide, M_L + 0.3, 3.05, BODY_W - 1.2, 1.6)
    p = tf.paragraphs[0]
    main = s.title.split("—")[0].strip()
    sub = s.title.split("—", 1)[1].strip() if "—" in s.title else ""
    r = p.add_run()
    r.text = main
    r.font.size, r.font.bold, r.font.name = Pt(44), True, FONT
    r.font.color.rgb = PAPER
    if sub:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(10)
        r2 = p2.add_run()
        r2.text = sub
        r2.font.size, r2.font.name = Pt(17), FONT
        r2.font.color.rgb = DARK_MUTED
    add_notes(slide, s.notes)
    return slide


def render_table(slide, b: Block, l, t, w, scale: float) -> float:
    rows, cols = len(b.rows), max(len(r) for r in b.rows)
    norm = [r + [""] * (cols - len(r)) for r in b.rows]
    hdr_h = 0.34 * scale
    row_h = 0.315 * scale
    shape = slide.shapes.add_table(rows, cols, Inches(l), Inches(t), Inches(w),
                                   Inches(hdr_h + row_h * (rows - 1)))
    table = shape.table
    table.first_row = True
    table.horz_banding = False

    # First column carries the labels and gets the extra room.
    if cols > 1:
        first = 0.46 if cols == 2 else 0.40
        table.columns[0].width = Emu(int(Inches(w) * first))
        rest = int((Inches(w) - table.columns[0].width) / (cols - 1))
        for c in range(1, cols):
            table.columns[c].width = Emu(rest)

    table.rows[0].height = Inches(hdr_h)
    for r in range(1, rows):
        table.rows[r].height = Inches(row_h)

    for r, row in enumerate(norm):
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.margin_left = Inches(0.09)
            cell.margin_right = Inches(0.07)
            cell.margin_top = Inches(0.035)
            cell.margin_bottom = Inches(0.035)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = WASH if r == 0 else PAPER
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            if r == 0:
                run = p.add_run()
                run.text = strip_marks(cell_text)
                run.font.size = Pt(TABLE_PT * scale)
                run.font.bold = True
                run.font.name = FONT
                run.font.color.rgb = INK
            else:
                add_runs(p, cell_text, TABLE_PT * scale, MUTED, bold_color=INK)
    return hdr_h + row_h * (rows - 1) + 0.20


def render_visual(slide, b: Block, l, t, w) -> float:
    h = 1.35
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = WASH
    box.line.color.rgb = ACCENT
    box.line.width = Pt(1.25)
    box.line.dash_style = 4  # dashed
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.25)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "IMAGE PLACEHOLDER"
    r.font.size, r.font.bold, r.font.name = Pt(9.5), True, FONT
    r.font.color.rgb = ACCENT
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)
    add_runs(p2, b.text, 12.0, MUTED)
    return h + 0.20


def render_blocks(slide, blocks: list[Block], l, t, w, scale: float):
    y = t
    for b in blocks:
        if b.kind == "table":
            y += render_table(slide, b, l, y, w, scale)
            continue
        if b.kind == "visual":
            y += render_visual(slide, b, l, y, w)
            continue

        if b.kind == "subhead":
            h = 0.30 * scale + 0.10
            _, tf = txbox(slide, l, y, w, h)
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = strip_marks(b.text).upper()
            r.font.size = Pt(SUBHEAD_PT * scale * 0.80)
            r.font.bold = True
            r.font.name = FONT
            r.font.color.rgb = ACCENT
            y += h
            continue

        if b.kind == "quote":
            size = QUOTE_PT * scale
            indent = 0.34
            lines = est_lines(b.text, size, w - indent)
            h = lines * (size * 1.30 / 72.0) + 0.16
            hline(slide, l, y + 0.03, 0.055, ACCENT, thick=h - 0.10)
            _, tf = txbox(slide, l + indent, y, w - indent, h)
            p = tf.paragraphs[0]
            add_runs(p, b.text, size, MUTED)
            for run in p.runs:
                run.font.italic = True
            y += h
            continue

        size = BODY_PT * scale
        indent = 0.30 * b.level
        bullet_pad = 0.26 if b.kind == "bullet" else 0.0
        lines = est_lines(b.text, size, w - indent - bullet_pad)
        h = lines * (size * 1.30 / 72.0) + (0.11 if b.kind in ("bullet", "number") else 0.16)

        if b.kind == "bullet":
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(l + indent + 0.02), Inches(y + size / 72.0 * 0.46),
                Inches(0.072), Inches(0.072))
            dot.fill.solid()
            dot.fill.fore_color.rgb = ACCENT
            dot.line.fill.background()
            dot.shadow.inherit = False

        _, tf = txbox(slide, l + indent + bullet_pad, y, w - indent - bullet_pad, h)
        p = tf.paragraphs[0]
        add_runs(p, b.text, size, MUTED, bold_color=INK)
        y += h
    return y


def split_columns(blocks: list[Block], w: float, scale: float):
    """Balance blocks across two columns, never splitting a table from its subhead."""
    groups: list[list[Block]] = []
    for b in blocks:
        if b.kind == "subhead" or not groups:
            groups.append([b])
        else:
            groups[-1].append(b)

    heights = [sum(block_height(b, w, scale) for b in g) for g in groups]
    total = sum(heights)
    left: list[Block] = []
    right: list[Block] = []
    acc = 0.0
    for g, h in zip(groups, heights):
        if acc + h / 2 <= total / 2 or not left:
            left.extend(g)
            acc += h
        else:
            right.extend(g)
    return left, right


SCALES = (1.0, 0.94, 0.88, 0.82)


def layout_plan(blocks: list[Block], avail_h: float):
    """Cheapest layout that fits every block, or None if they will not fit at all."""
    colw = (BODY_W - COL_GAP) / 2
    for scale in SCALES:
        if sum(block_height(b, BODY_W, scale) for b in blocks) <= avail_h:
            return "one", scale
        lft, rgt = split_columns(blocks, colw, scale)
        h2 = max(sum(block_height(b, colw, scale) for b in lft),
                 sum(block_height(b, colw, scale) for b in rgt))
        if h2 <= avail_h:
            return "two", scale
    return None


def max_prefix(blocks: list[Block], avail_h: float, scale: float) -> int:
    """Largest leading run of blocks that still fits two columns at `scale`."""
    colw = (BODY_W - COL_GAP) / 2
    for n in range(len(blocks), 0, -1):
        lft, rgt = split_columns(blocks[:n], colw, scale)
        h = max(sum(block_height(b, colw, scale) for b in lft),
                sum(block_height(b, colw, scale) for b in rgt))
        if h <= avail_h:
            return n
    return 1


def _open_content_slide(prs, title: str, headline: str):
    """Draw chrome (label, title, rule, headline); return (slide, body_top, body_height)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, PAPER)

    label, _, rest = title.partition("—")
    label, rest = label.strip(), rest.strip()
    _, tf = txbox(slide, M_L, M_T, BODY_W, 0.52)
    p = tf.paragraphs[0]
    if rest and re.match(r"^Slide\s+\d+$", label):
        r0 = p.add_run()
        r0.text = label.upper() + "   "
        r0.font.size, r0.font.bold, r0.font.name = Pt(11), True, FONT
        r0.font.color.rgb = ACCENT
        title_text = rest
    else:
        title_text = title
    r = p.add_run()
    r.text = title_text
    r.font.size, r.font.bold, r.font.name = Pt(25), True, FONT
    r.font.color.rgb = INK

    y = M_T + 0.58
    hline(slide, M_L, y, BODY_W)
    y += 0.20

    if headline:
        size = 17.0
        lines = est_lines(headline, size, BODY_W)
        h = lines * (size * 1.32 / 72.0) + 0.10
        _, tfh = txbox(slide, M_L, y, BODY_W, h)
        ph = tfh.paragraphs[0]
        add_runs(ph, headline, size, INK, bold_color=ACCENT)
        for run in ph.runs:
            run.font.bold = True
        y += h + 0.22

    return slide, y, SLIDE_H - M_B - y


def _place(slide, blocks, mode, scale, y):
    if mode == "one":
        render_blocks(slide, blocks, M_L, y, BODY_W, scale)
    else:
        colw = (BODY_W - COL_GAP) / 2
        lft, rgt = split_columns(blocks, colw, scale)
        render_blocks(slide, lft, M_L, y, colw, scale)
        render_blocks(slide, rgt, M_L + colw + COL_GAP, y, colw, scale)


def render_content(prs, s: Slide, warn: list[str]):
    """Render one markdown section, spilling onto '(cont.)' slides when it overruns."""
    remaining = list(s.blocks)
    first = True

    while True:
        title = s.title if first else f"{s.title}  (cont.)"
        slide, y, avail = _open_content_slide(prs, title, s.headline if first else "")

        if remaining:
            plan = layout_plan(remaining, avail)
            if plan:
                _place(slide, remaining, plan[0], plan[1], y)
                remaining = []
            else:
                scale = SCALES[-1]
                n = max_prefix(remaining, avail, scale)
                _place(slide, remaining[:n], "two", scale, y)
                remaining = remaining[n:]
                if n == 1 and block_height(remaining[0] if remaining else s.blocks[0],
                                           (BODY_W - COL_GAP) / 2, scale) > avail:
                    warn.append(f"{s.title} (single block taller than one slide)")

        if first:
            add_notes(slide, s.notes)
            first = False
        if not remaining:
            return slide


def add_notes(slide, notes: list[str]):
    if not notes:
        return
    tf = slide.notes_slide.notes_text_frame
    tf.text = ""
    first = True
    for n in notes:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = strip_marks(n)
        r.font.size = Pt(11)


# ---------------------------------------------------------------- main

def build(src: Path, out: Path) -> None:
    slides = parse(src.read_text(encoding="utf-8"))

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    warn: list[str] = []
    for s in slides:
        if s.kind == "title":
            render_title(prs, s)
        elif s.kind == "divider":
            render_divider(prs, s)
        else:
            render_content(prs, s, warn)

    prs.save(out)

    print(f"wrote {out}  ({len(prs.slides)} slides)")
    noted = sum(1 for s in slides if s.notes)
    print(f"speaker notes attached to {noted} slides")
    if warn:
        print("\nDENSE - may overflow, consider splitting:")
        for t in warn:
            print(f"  - {t}")


if __name__ == "__main__":
    src = Path(sys.argv[1]) if len(sys.argv) > 1 else Path(__file__).parent / "qlaudia_pitch_deck.md"
    out = Path(sys.argv[2]) if len(sys.argv) > 2 else src.with_suffix(".pptx")
    build(src, out)
